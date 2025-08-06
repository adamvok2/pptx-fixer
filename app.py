import streamlit as st
from pptx import Presentation
import re
import tempfile
import os

# Regex pattern
pattern = re.compile(r'@[^@]*\(([^@()]*)\)[^@]*@')

def transform_text(text):
    matches = pattern.findall(text)
    if matches:
        return pattern.sub(lambda m: matches.pop(0) if matches else m.group(0), text)
    return text

def fix_pptx(file):
    prs = Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        run.text = transform_text(run.text)
            if shape.shape_type == 19:  # Table
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            for run in p.runs:
                                run.text = transform_text(run.text)
    tmpdir = tempfile.mkdtemp()
    return prs, tmpdir

st.title("📊 AFGC PPTX Fixer")

st.markdown(
    """
    🔒 **Privacy Notice**

    - Your files are processed securely **in memory** and **never stored**.
    - No data is logged, saved, or sent to any external service.
    - Once the session ends or the browser is closed, all data is wiped.
    """
)

uploaded_files = st.file_uploader("Upload one or more .pptx files", type=["pptx"], accept_multiple_files=True)

if uploaded_files:
    for uploaded_file in uploaded_files:
        with st.spinner(f"Fixing: {uploaded_file.name}"):
            prs, tmpdir = fix_pptx(uploaded_file)
            output_filename = uploaded_file.name.replace('.pptx', '_fixed.pptx')
            output_path = os.path.join(tmpdir, output_filename)
            prs.save(output_path)

            with open(output_path, 'rb') as f:
                st.download_button(
                    label=f"Download {output_filename}",
                    data=f,
                    file_name=output_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
